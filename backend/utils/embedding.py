from sentence_transformers import SentenceTransformer
import numpy as np
import fitz  # PyMuPDF
import json
import os
import shutil
import mysql.connector
from docx import Document
import pandas as pd
from pptx import Presentation
from bs4 import BeautifulSoup
import zipfile

# Khởi tạo mô hình embedding
model = SentenceTransformer('all-MiniLM-L6-v2')

# Chunk văn bản thành từng phần nhỏ
def chunk_text(text, size=500):
    return [text[i:i + size] for i in range(0, len(text), size)]

# Kết nối MySQL nếu cần
def get_db():
    return mysql.connector.connect(
        host="localhost",
        port=3300,
        user="root",
        password="",
        database="chatbot"
    )

# Trích xuất văn bản từ file
def extract_text_from_file(filepath):
    ext = filepath.split('.')[-1].lower()
    if ext == 'pdf':
        doc = fitz.open(filepath)
        return "\n".join([page.get_text() for page in doc])
    elif ext == 'docx':
        doc = Document(filepath)
        return "\n".join([para.text for para in doc.paragraphs])
    elif ext in ['txt', 'md']:
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()
    elif ext == 'csv':
        df = pd.read_csv(filepath)
        return df.astype(str).to_string(index=False)
    elif ext == 'xlsx':
        df = pd.read_excel(filepath, sheet_name=None)
        return "\n".join(df[sheet].astype(str).to_string(index=False) for sheet in df)
    elif ext == 'pptx':
        prs = Presentation(filepath)
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    elif ext == 'html':
        with open(filepath, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
            return soup.get_text()
    elif ext == 'zip':
        text_content = []
        temp_dir = 'temp_zip_extract'
        os.makedirs(temp_dir, exist_ok=True)
        with zipfile.ZipFile(filepath, 'r') as zipf:
            for name in zipf.namelist():
                if name.endswith(('.txt', '.md', '.csv', '.docx', '.pdf', '.html', '.xlsx', '.pptx')):
                    zipf.extract(name, temp_dir)
                    full_path = os.path.join(temp_dir, name)
                    try:
                        text = extract_text_from_file(full_path)
                        text_content.append(f"### {name}\n" + text)
                    except Exception as e:
                        text_content.append(f"[Bỏ qua file {name}]: {str(e)}")
        shutil.rmtree(temp_dir)
        return "\n\n".join(text_content)
    else:
        raise Exception("Unsupported format (chỉ hỗ trợ PDF, DOCX, TXT, MD, CSV, XLSX, PPTX, HTML, ZIP)")

# Xử lý và lưu embedding
def process_and_store(filepath, doc_id, db=None):
    text = extract_text_from_file(filepath)
    chunks = chunk_text(text)
    vectors = model.encode(chunks)

    close_db = False
    if db is None:
        db = get_db()
        close_db = True

    cursor = db.cursor()
    for chunk, vector in zip(chunks, vectors):
        cursor.execute(
            "INSERT INTO embeddings (document_id, chunk, vector) VALUES (%s, %s, %s)",
            (doc_id, chunk, json.dumps(vector.tolist()))
        )
    db.commit()

    if close_db:
        db.close()
